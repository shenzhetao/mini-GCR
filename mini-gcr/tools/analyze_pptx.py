# -*- coding: utf-8 -*-
import sys
import os
os.environ['PYTHONIOENCODING'] = 'utf-8'

from pptx import Presentation

# Open the PPTX file
pptx_path = r'c:\project\mini-GCR\第2次作业-12-基于生成式推荐的电商支付界面与购物车互补品推荐算法研究.pptx'

prs = Presentation(pptx_path)

output_lines = []

# Total number of slides
total_slides = len(prs.slides)
output_lines.append(f'Total number of slides: {total_slides}')
output_lines.append('')

# Slide layout names
output_lines.append('Slide Layout Names Used:')
for i, slide_layout in enumerate(prs.slide_layouts):
    output_lines.append(f'  Layout {i}: {slide_layout.name}')
output_lines.append('')

# Details for each slide
output_lines.append('='*60)
output_lines.append('SLIDE DETAILS')
output_lines.append('='*60)

for idx, slide in enumerate(prs.slides, 1):
    output_lines.append('')
    output_lines.append(f'Slide {idx}:')
    
    # Get slide title
    title = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                title = text[:100]
                break
    
    if title:
        output_lines.append(f'  Title: {title}')
    else:
        output_lines.append(f'  Title: (No title)')
    
    # Number of shapes
    num_shapes = len(slide.shapes)
    output_lines.append(f'  Number of shapes: {num_shapes}')

output_text = '\n'.join(output_lines)

# Print to console
print(output_text)

# Save to file
output_path = r'c:\project\mini-GCR\mini-gcr\pptx_analysis.txt'
with open(output_path, 'w', encoding='utf-8') as f:
    f.write(output_text)

print(f'\n[Output saved to: {output_path}]')
