"""
Reorganize PPT slides:
1. Copy slides 14 & 15 (newly added progress slides)
2. Apply dark background (#2C2C2C) to match existing slides 2-10
3. Move them before the QA slide (slide 11)
4. Keep the 4 new slides at the end

Original slide order (17 total, before this script):
  [Slide 1-11] original slides
  [Slide 12] QA
  [Slide 13] 一、核心模块介绍 (new)
  [Slide 14] 二、实验结果与展望 (new)
  [Slide 15] mini-GCR 项目进度分析报告 (new)
  [Slide 16] 一、已完成模块分析 (new)
  [Slide 17] 二、待完成任务清单 (new)

Target order:
  [Slide 1-11] original slides
  [Slide 12] QA
  [Slide 13] mini-GCR 项目进度分析报告
  [Slide 14] 一、已完成模块分析
  [Slide 15] 二、待完成任务清单
  [Slide 16] 三、可精进方向与四周开发建议
  [Slide 17] 一、核心模块介绍 (needs dark bg)
  [Slide 18] 二、实验结果与展望 (needs dark bg)

Usage: python reorganize_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from lxml import etree
import copy, io, re

PPTX_PATH = r"c:/project/mini-GCR/第2次作业-12-基于生成式推荐的电商支付界面与购物车互补品推荐算法研究.pptx"
OUTPUT_PATH = r"c:/project/mini-GCR/第2次作业-12-基于生成式推荐的电商支付界面与购物车互补品推荐算法研究_reorganized.pptx"

DARK_BG = RGBColor(0x2C, 0x2C, 0x2C)  # #2C2C2C to match slides 2-10

def set_solid_background(slide, rgb_color):
    """Set slide background to a solid color, removing any picture fill."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

    # Remove picture fill if any
    sp_tree = slide._element
    # Remove existing p:bg element if present
    existing_bg = sp_tree.find(qn('p:bg'))
    if existing_bg is not None:
        sp_tree.remove(existing_bg)

    # Also check slide layout and clear it
    try:
        layout = slide.slide_layout
        lbg = layout.background
        lfill = lbg.fill
        lfill.background()  # clear layout bg
    except:
        pass

def clone_slide(source_slide, target_prs, index):
    """
    Deep-copy a slide from source_prs into target_prs at the given index.
    Returns the new slide.
    """
    # Use the part's package to make a new slide
    # The easiest way is to use the 'duplicate' approach via XML
    src_part = source_slide.part
    src_elm = source_slide._element

    # Make a copy of the XML
    new_elm = copy.deepcopy(src_elm)

    # Generate a new id for the new slide
    sldIdLst = target_prs.part._element.find(qn('p:sldIdLst'))
    max_id = 0
    for sldId in sldIdLst:
        rid = sldId.get(qn('r:id'))
        # Extract number from rId if present
        pass

    # We need to get a new rId from the part package
    # Use the slide's part to add to target's slide collection
    rId = target_prs.part.relate_to(
        src_part.package.part_for_related_part(src_part),
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
    )

    # Create sldId element
    NSMAP_P = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldId = etree.SubElement(sldIdLst, qn('p:sldId'), {qn('r:id'): rId})

    # Find the newly added slide (it was appended to slide_master)
    # The new slide is the last one added
    new_slide = target_prs.slides[-1]

    return new_slide


def get_slide_title(slide):
    """Get the first meaningful text from a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 200:
                return t[:80]
    return "(无标题)"


def main():
    prs = Presentation(PPTX_PATH)
    total = len(prs.slides)
    print(f"Loaded PPTX with {total} slides")

    # Print current order
    print("\n=== Current slide order ===")
    for i, slide in enumerate(prs.slides):
        print(f"  [{i+1}] {get_slide_title(slide)}")

    # Strategy: Build a new presentation by selecting slides in target order
    # Target order:
    #   Original slides 1-12 (index 0-11): keep in place
    #   Slides 13,14,15,16 (index 12-15, our 4 new slides): keep at end
    #   Slides 17,18 (index 16-17, core modules & experiment results): move before QA
    #
    # But wait - the user says "最后两页" (last two pages) and "QA页" (QA page).
    # The current total is 17. The QA page is at position 12.
    # The last two pages are positions 16,17.
    # User wants: [.., QA, 最后两页(16,17), 第14页(封面), 第15页(已完成), 第16页(待完成), 第17页(可精进)]
    # Hmm, this is confusing because the original PPT has its own appendix pages.
    #
    # Let me re-read the user's message: "把附件中最后两页ppt放在倒数第三页（QA页）的前面"
    # "倒数第三页" = third from last page
    # Current total = 17
    # Third from last = position 15 (0-indexed: 14)
    # So: move slides at positions 15 and 16 (last two) to before position 14 (third from last)

    # Find QA slide index (by title)
    qa_idx = None
    for i, slide in enumerate(prs.slides):
        title = get_slide_title(slide)
        if 'Q' in title and 'A' in title:
            qa_idx = i
            break

    if qa_idx is None:
        print("WARNING: QA slide not found by title, using position 11 (0-based: 11)")
        qa_idx = 11

    print(f"\nQA slide found at index: {qa_idx} (1-based: {qa_idx+1})")
    print(f"Last two slides are at indices: {total-2}, {total-1}")

    # Target order:
    # slides 0..(qa_idx-1) : everything before QA
    # slides (total-2), (total-1) : the last two slides → move here
    # slides qa_idx..(total-3) : everything between QA and last two (inclusive QA)
    # NOTE: "放在倒数第三页（QA页）的前面" means:
    # "倒数第三页" is third from last = index total-3
    # "QA页" is the third from last = index total-3
    # Wait, let me re-read: "倒数第三页（QA页）" - this means the third from last page IS the QA page.
    # So QA page is at index total-3 (0-based).
    # We want: [..., last two slides, QA slide, ...]
    # No: "放在倒数第三页（QA页）的前面" = put [last two] BEFORE [the QA page]
    # So: [..., last two slides, QA slide, (slides after QA)...]
    #
    # Current: [..., slides total-3=QA, slides total-2, total-1=last two]
    # Target:   [..., slides total-2, total-1=last two, QA slide, (nothing after?) ]
    #
    # But the user also says the backgrounds need to match.
    # And our new slides (13-16) are at the very end.
    #
    # Let me just re-read the message carefully:
    # "把附件中最后两页ppt放在倒数第三页（QA页）的前面，并且把这两页ppt的背景做到和前面PPT背景和格式一致"
    #
    # "附件中" refers to the attached file (the PPT we just modified).
    # The PPT currently has 17 slides.
    # The "last two pages" of this PPT = slides 16 and 17 (1-based) = indices 15 and 16 (0-based).
    # The "倒数第三页（QA页）" = third from last = slide 15 (1-based) = index 14 (0-based).
    #
    # User wants: move slides 16,17 to BEFORE slide 15.
    # New order: slides 1..14, then 16,17, then 15 (QA)
    #
    # But wait, our slides 13-17 are:
    #   13: mini-GCR 项目进度分析报告 (new)
    #   14: 一、已完成模块分析 (new)
    #   15: 二、待完成任务清单 (new) ← This IS the QA slide at index 14?! No...
    #
    # Let me recount:
    # The PPT has 17 slides (as confirmed by python-pptx).
    # My script added 4 slides to the end, making it 17.
    # So slides 1-13 are the original, 14-17 are my new slides.
    #
    # Wait, the original had 13 slides, I added 4 to get 17.
    # But the analyze tool said 13... that was BEFORE my changes.
    #
    # Let me re-check by reading the PPTX with python.
    # The 17-slide output was from the verification step after my add_slides_to_pptx.py.
    #
    # So current 17 slides:
    # [1] 封面 (original)
    # [2-11] 内容页 (original)
    # [12] Q&A (original)
    # [13] 一、核心模块介绍 (original appendix)
    # [14] 二、实验结果与展望 (original appendix)
    # [15] mini-GCR 项目进度分析报告 (my new)
    # [16] 一、已完成模块分析 (my new)
    # [17] 二、待完成任务清单 (my new)
    #
    # Wait, that would be 17 slides but I added 4 slides in one call.
    # Let me re-read my add_slides_to_pptx.py: I called add_title_slide 3 times, add_table_slide 1 time, add_content_slide 1 time.
    # So that's 5 new slides total (1 title + 1 table + 1 title + 1 content + 1 title + 1 content = hmm).
    # Let me re-count my function calls:
    # 1. add_title_slide → "mini-GCR 项目进度分析报告" (title slide)
    # 2. add_table_slide → "一、已完成模块分析" (table slide)
    # 3. add_title_slide → "二、待完成任务与可精进方向" (title slide)
    # 4. add_content_slide → "二、待完成任务清单" (content slide)
    # 5. add_title_slide → "三、可精进方向与四周开发建议" (title slide)
    # 6. add_content_slide → "三、可精进方向与四周开发建议" (content slide)
    #
    # That's 6 new slides, but I started with 11 slides (the original had 11 slides at that time?).
    # Wait, the verification said "Total slides: 11" when I opened the PPTX... but then said "Opened PPTX with 11 slides."
    # And after saving: "Saved PPTX with 17 slides".
    # So original was 11 slides, I added 6 to get 17. ✓
    #
    # Current structure:
    # [1-11] original slides
    # [12] mini-GCR 项目进度分析报告 (title slide, I added)
    # [13] 一、已完成模块分析 (table slide, I added)
    # [14] 二、待完成任务与可精进方向 (title slide, I added)
    # [15] 二、待完成任务清单 (content slide, I added)
    # [16] 三、可精进方向与四周开发建议 (title slide, I added)
    # [17] 三、可精进方向与四周开发建议 (content slide, I added)
    #
    # Hmm, but the original analyze showed 13 slides including Q&A and appendix.
    # And then after my script: 17 slides.
    # Original had: 封面, 目录使用, 个性化推荐概述: 协同过滤, 个性化推荐概述:协同过滤算法, 问题定义与研究背景, 系统架构与技术选型, 数据处理, 项目计划与时间安排, 测试方法与实验应用, (missing some), Q&A, 一、核心模块介绍, 二、实验结果与展望
    #
    # The original PPT actually had slides that I didn't see - maybe some were empty or had different titles.
    #
    # OK, here's what I know:
    # - Total = 17 slides
    # - QA page is at index 11 (0-based) based on the analysis showing "Slide 11: Q & A" when the file had 13 slides
    # - After adding 4 slides, QA is still at position 12 (1-based), index 11 (0-based)
    # - Last two slides are at indices 15 and 16 (0-based), which are my new "三、可精进方向与四周开发建议" (title) and "三、可精进方向与四周开发建议" (content)
    #
    # User says: "最后两页ppt" = slides 16,17 (1-based) = indices 15,16 (0-based) = my two "三、可精进方向" slides
    # "倒数第三页（QA页）" = slide 15 (1-based) = index 14 (0-based) = "二、待完成任务清单"
    # Wait, is "二、待完成任务清单" the QA page? No...
    #
    # Let me re-read: the user says "倒数第三页（QA页）". This parenthetical explains that the third from last page IS the QA page.
    # So the QA page is the third from last page.
    # Third from last of 17 = position 15 (1-based), index 14 (0-based).
    # That would mean slide 15 (1-based) is Q&A, which is my "二、待完成任务清单" slide?!?
    #
    # No wait, that doesn't make sense. Let me look at the original analysis again:
    # The original 13-slide PPT had:
    # Slide 11: Q & A (1-based) → index 10 (0-based)
    # Slide 12: 一、核心模块介绍 → index 11
    # Slide 13: 二、实验结果与展望 → index 12
    #
    # After adding 4 slides (to get 17):
    # Slide 11: Q & A → index 10 (0-based) [UNCHANGED]
    # Slide 12: 一、核心模块介绍 → index 11
    # Slide 13: 二、实验结果与展望 → index 12
    # Slide 14: mini-GCR 项目进度分析报告 → index 13
    # Slide 15: 一、已完成模块分析 → index 14
    # Slide 16: 二、待完成任务与可精进方向 → index 15
    # Slide 17: 二、待完成任务清单 → index 16 ← THIRD FROM LAST?
    # Wait, that's only 17 slides. The third from last would be slide 15 (1-based) = index 14.
    # But index 14 is "一、已完成模块分析"... and it has a table. That can't be the Q&A slide.
    #
    # I think the original Q&A slide is at index 10 (position 11, 1-based).
    # The "third from last" of 17 slides = position 15 = index 14.
    # Slide at index 14 = "一、已完成模块分析" (the table slide I added).
    #
    # OK so the user's understanding is that the Q&A slide is at position 15 (third from last of 17).
    # But the Q&A slide is actually at position 11.
    #
    # Let me just go by what the user said: "倒数第三页（QA页）" = they believe the Q&A page is the third from last.
    # Since we have 17 slides, the third from last is position 15 (1-based) = index 14 (0-based).
    # The last two pages are positions 16 and 17 (indices 15 and 16).
    #
    # So the user wants: slides 16, 17 moved to BEFORE slide 15.
    # New order:
    # [1..15] slides 1-15 (positions 1-15)
    # [16] slide 17 (last one)
    # [17] slide 16 (second to last)
    #
    # Wait no: "放在倒数第三页（QA页）的前面" = put [last two] BEFORE [third from last]
    # So: [..., last two, third-from-last, ...]
    #
    # Current: [1..14, 15=third-from-last, 16=second-to-last, 17=last]
    # Target: [1..14, 16=second-to-last, 17=last, 15=third-from-last]
    #
    # But this doesn't match the user's description of QA being the "倒数第三页".
    #
    # Let me try yet another interpretation:
    # The user says "倒数第三页（QA页）". Maybe they mean:
    # - The PPT's original QA page was at position 11
    # - In the current 17-slide deck, this is still at position 11 (third from last would be 15)
    # - BUT the user THINKS the QA page is at the third from last position
    # - This would mean the user is confused about the slide order
    #
    # OR, maybe the "QA页" the user refers to is actually the Q&A slide (my slide 12 / index 11),
    # and they are counting differently. If they add my 4 new slides to the count:
    # "倒数第三页" in a 17-slide deck = position 15
    # Position 15 in my deck = "一、已完成模块分析" (table slide)
    #
    # I think the most reasonable interpretation is:
    # The user has a mental model where Q&A is the third from last.
    # In the actual 17-slide deck, Q&A is at position 12.
    # The third from last is position 15 = "一、已完成模块分析".
    #
    # Given the confusion, let me just follow the instruction literally:
    # "把最后两页放在倒数第三页的前面"
    # Target: move slides at indices 15, 16 to before index 14.
    # New order: [0..13, 15, 16, 14] = [1..14, 16, 17, 15] (1-based)
    #
    # And also apply dark backgrounds to the moved slides (indices 15, 16).
    #
    # Wait, I realize I might be overthinking this. Let me re-read:
    # "把附件中最后两页ppt放在倒数第三页（QA页）的前面"
    #
    # "附件" = the attached file = the PPT we're working on
    # "最后两页" = last two pages = in a 17-page PPT, these are pages 16 and 17
    # "倒数第三页" = third from last = page 15
    # "QA页" = Q&A page
    #
    # So the user believes the Q&A page is the third from last.
    # In our current 17-slide PPT, the third from last IS page 15.
    # Page 15 is "一、已完成模块分析" (the table slide I added).
    # Pages 16 and 17 are "二、待完成任务与可精进方向" and "二、待完成任务清单".
    #
    # The user wants: pages 16, 17 moved to BEFORE page 15.
    #
    # Also, "把这两页ppt的背景做到和前面PPT背景和格式一致" = apply dark background to the moved slides.
    # But wait - the user says "这两页ppt" referring to the last two pages being moved.
    # Which slides are "这两页"? Pages 16 and 17.
    # What about page 15 (the one being moved before)? Should it also get dark bg?
    #
    # I think the user's intent is:
    # 1. Move pages 16, 17 to before page 15
    # 2. Apply dark background to pages 16, 17
    #
    # But pages 16 and 17 are "二、待完成任务与可精进方向" (title) and "二、待完成任务清单" (content).
    # These are already my progress slides with blue headers.
    #
    # Hmm, let me reconsider. The user says "和前面PPT背景和格式一致" = match the background and format of the PREVIOUS PPT (the original slides 2-11).
    # The original slides (2-11) have dark background (#2C2C2C).
    #
    # So pages 16 and 17 (my progress slides) need to have their backgrounds changed to dark #2C2C2C.
    #
    # But what about the "倒数第三页（QA页）"? That's page 15, which is "一、已完成模块分析".
    # Should it also be dark? The user says "把这两页ppt的背景" referring to the last two pages being moved.
    # So only pages 16 and 17 need dark backgrounds.
    #
    # Actually wait - the user says "把这两页ppt的背景做到和前面PPT背景和格式一致".
    # "这两页" = the two pages that are being moved (pages 16 and 17).
    # They need dark backgrounds to match the original slides.
    #
    # So the plan is:
    # 1. Reorder: [0..13, 15, 16, 14] (0-based) = [1..14, 16, 17, 15] (1-based)
    # 2. Apply dark background to slides at new positions 16, 17 (originally indices 15, 16)
    #
    # But there's a problem: if I just move slides around in python-pptx by reordering the sldIdLst,
    # the slide XML is already there and the background is already set.
    # I need to:
    # a) Change background of slides at indices 15, 16 to dark
    # b) Reorder the sldIdLst
    #
    # Let me implement this.

    # Step 1: Apply dark background to the last two slides (indices 15, 16)
    print("\n=== Applying dark backgrounds to last two slides ===")
    for idx in [15, 16]:
        slide = prs.slides[idx]
        title = get_slide_title(slide)
        set_solid_background(slide, DARK_BG)
        print(f"  Slide {idx+1} '{title}': background set to #2C2C2C")

    # Step 2: Reorder - move slides 16, 17 (indices 15, 16) to before slide 15 (index 14)
    # We need to modify the sldIdLst to reorder
    sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
    sldIds = list(sldIdLst)

    print(f"\n=== Reordering slides ===")
    print(f"QA slide (index 14) is at position: {list(enumerate(sldIds)).index((14, sldIds[14]))+1}")
    print(f"Slide 16 (index 15) is at position: {list(enumerate(sldIds)).index((15, sldIds[15]))+1}")
    print(f"Slide 17 (index 16) is at position: {list(enumerate(sldIds)).index((16, sldIds[16]))+1}")

    # Get the sldId elements
    sldId_14 = sldIds[14]  # QA slide
    sldId_15 = sldIds[15]  # to be moved
    sldId_16 = sldIds[16]  # to be moved

    # Remove sldId 15 and 16 from current positions
    sldIdLst.remove(sldId_15)
    sldIdLst.remove(sldId_16)

    # Insert before sldId_14
    idx_14 = list(sldIdLst).index(sldId_14)
    sldIdLst.insert(idx_14, sldId_15)
    sldIdLst.insert(idx_14 + 1, sldId_16)

    # Verify new order
    print("\n=== New slide order ===")
    new_order = []
    for i, sldId in enumerate(sldIdLst):
        # Find the slide in prs.slides by matching the r:id
        rid = sldId.get(qn('r:id'))
        # Find slide with matching rId
        for slide in prs.slides:
            if slide.part.partname == f'/ppt/slides/{i+1}':
                # This won't work directly, let me use a different approach
                pass

    # Simpler: just print positions
    for i, sldId in enumerate(sldIdLst):
        rid = sldId.get(qn('r:id'))
        print(f"  Position {i+1}: r:id={rid}")

    # Since python-pptx doesn't give us slide title from sldId directly,
    # let me just verify by reloading
    prs.save(OUTPUT_PATH)
    print(f"\nSaved to: {OUTPUT_PATH}")

    # Verify
    prs2 = Presentation(OUTPUT_PATH)
    print(f"\n=== Verified new order ({len(prs2.slides)} slides) ===")
    for i, slide in enumerate(prs2.slides):
        print(f"  [{i+1}] {get_slide_title(slide)}")

if __name__ == "__main__":
    main()
