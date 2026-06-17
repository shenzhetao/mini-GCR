"""
Analyze PPTX structure and backgrounds.
Usage: python analyze_pptx_style.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy, re

PPTX_PATH = r"c:/project/mini-GCR/第2次作业-12-基于生成式推荐的电商支付界面与购物车互补品推荐算法研究.pptx"

def get_bg_info(slide):
    """Extract background fill info from a slide."""
    bg = slide.background
    fill = bg.fill
    info = {"type": fill.type.name if fill else "None"}

    # Check for solid fill
    if fill.type == 1:  # SOLID
        info["solid_fg"] = fill.fore_color.rgb if fill.fore_color else None
        info["solid_theme"] = fill.fore_color.theme_color if hasattr(fill.fore_color, 'theme_color') else None

    # Check slide layout background
    layout = slide.slide_layout
    lfill = layout.background.fill
    info["layout_fill"] = lfill.type.name if lfill else "None"
    if lfill.type == 1:
        try:
            info["layout_solid"] = lfill.fore_color.rgb
        except:
            info["layout_solid"] = "theme_color"

    # Check for picture fill
    if fill.type == 8:  # PICTURE
        info["has_picture"] = True

    # Get slide XML to check for shape backgrounds
    slide_xml = slide._element
    # Look for bg element
    bg_elem = slide_xml.find(qn('p:bg'))
    if bg_elem is not None:
        info["bg_element"] = etree.tostring(bg_elem, pretty_print=True).decode()[:500]

    return info

def describe_slide(slide, idx):
    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 200:
                title = t[:80]
                break
    bg_info = get_bg_info(slide)
    return f"Slide {idx}: '{title}' | BG: {bg_info}"

def main():
    prs = Presentation(PPTX_PATH)
    print(f"Total slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        info = describe_slide(slide, i+1)
        print(info)

    # Print slide dimensions
    print(f"\nSlide size: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

if __name__ == "__main__":
    main()
