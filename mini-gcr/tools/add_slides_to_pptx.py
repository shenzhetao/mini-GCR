"""
Add 3 progress analysis slides to the existing PPTX.
Usage: python add_slides_to_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# Paths
PPTX_PATH = r"c:/project/mini-GCR/第2次作业-12-基于生成式推荐的电商支付界面与购物车互补品推荐算法研究.pptx"
OUTPUT_PATH = r"c:/project/mini-GCR/第2次作业-12-基于生成式推荐的电商支付界面与购物车互补品推荐算法研究_updated.pptx"

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

# Colors matching the project's theme (blue-ish academic)
TITLE_COLOR = hex_to_rgb("1F4E79")   # Dark blue
ACCENT_COLOR = hex_to_rgb("2E75B6")  # Medium blue
GREEN_COLOR = hex_to_rgb("375623")   # Dark green for done
RED_COLOR = hex_to_rgb("C00000")     # Red for issues
LIGHT_BG = hex_to_rgb("DEEAF6")      # Light blue background for boxes
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = hex_to_rgb("595959")
LIGHT_GRAY = hex_to_rgb("F2F2F2")

def add_title_slide(prs, title_text, subtitle_text=""):
    """Add a section divider / title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Full-width title bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_COLOR
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(12), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle_text:
        subBox = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(12), Inches(0.6))
        stf = subBox.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(18)
        sp.font.color.rgb = GRAY
        sp.font.italic = True
        sp.alignment = PP_ALIGN.LEFT

    return slide

def add_content_slide(prs, title, bullets, footer_text=""):
    """Add a bullet-point content slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_COLOR
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.18), Inches(12), Inches(0.75))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area
    content_top = Inches(1.2)
    content_height = Inches(5.3)
    content_width = prs.slide_width - Inches(0.6)

    txBox2 = slide.shapes.add_textbox(Inches(0.3), content_top, content_width, content_height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        # Handle sub-bullets
        indent_level = bullet.get("level", 0)
        indent_inches = 0.3 + indent_level * 0.35

        p.text = bullet.get("text", "")
        p.level = indent_level

        run = p.runs[0] if p.runs else p.add_run()
        run.text = p.text
        p.runs[0].font.size = Pt(16 - indent_level * 1.5)
        p.runs[0].font.color.rgb = GRAY if indent_level > 0 else BLACK

        if bullet.get("bold"):
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = ACCENT_COLOR

        p.space_before = Pt(4)
        p.space_after = Pt(2)

    # Footer
    if footer_text:
        ftBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(12), Inches(0.4))
        ftf = ftBox.text_frame
        fp = ftf.paragraphs[0]
        fp.text = footer_text
        fp.font.size = Pt(10)
        fp.font.color.rgb = GRAY
        fp.font.italic = True

    return slide

def add_table_slide(prs, title, headers, rows, footer_text=""):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_COLOR
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.18), Inches(12), Inches(0.75))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_width = prs.slide_width - Inches(0.6)
    col_width = table_width / num_cols

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(1.2),
        table_width, Inches(num_rows * 0.45)
    ).table

    # Style header
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(13)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # Style rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = BLACK
            para.alignment = PP_ALIGN.LEFT

    # Footer
    if footer_text:
        ftBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(12), Inches(0.4))
        ftf = ftBox.text_frame
        fp = ftf.paragraphs[0]
        fp.text = footer_text
        fp.font.size = Pt(10)
        fp.font.color.rgb = GRAY
        fp.font.italic = True

    return slide


def main():
    prs = Presentation(PPTX_PATH)
    print(f"Opened PPTX with {len(prs.slides)} slides.")

    # ─── Slide 1: 项目完成情况总览 ───────────────────────────────────────────
    add_title_slide(
        prs,
        "mini-GCR 项目进度分析报告",
        "生成式互补品推荐系统 | 分析日期：2026-05-23"
    )

    headers1 = ["模块", "文件", "状态", "完成度", "说明"]
    rows1 = [
        ["数据管线", "generate_mock_data.py / preprocess.py", "✅ 代码完成", "80%", "数据目录为空，需先运行流水线"],
        ["SASRec 基线", "models/sasrec/model.py + train_sasrec.py", "✅ 代码完成", "90%", "未训练，checkpoints 为空"],
        ["minGPT 生成模型", "models/mingpt/model.py + trainer.py + decode.py", "✅ 代码完成", "85%", "存在 beam search 多样性 bug"],
        ["推理服务", "services/inference.py + reason.py", "✅ 代码完成", "85%", "静默降级，理由模板仅 5 个类目对"],
        ["评估脚本", "scripts/eval.py", "✅ 代码完成", "90%", "硬编码评测上限 100 条，未执行"],
        ["FastAPI", "app.py", "✅ 代码完成", "90%", "4 个接口正常，含 CORS"],
        ["Vue 前端", "frontend/src/App.vue", "✅ 基础完成", "70%", "缺少对比视图、指标仪表盘、商品图片"],
        ["基础设施", "config.py + run_pipeline.py + requirements.txt", "✅ 完整", "95%", "配置中心、流水线脚本、依赖管理"],
    ]
    add_table_slide(
        prs,
        "一、已完成模块分析",
        headers1, rows1,
        "注：整体代码完成度约 80%，但数据目录为空、模型未训练，端到端链路尚未打通。"
    )

    # ─── Slide 2: 待完成任务 ────────────────────────────────────────────────
    add_title_slide(prs, "二、待完成任务与可精进方向", "P0 阻塞项 | P1/P2 优化方向")

    bullets2 = [
        {"text": "■ P0 必须完成（影响核心交付）", "bold": True},
        {"text": "① 生成 Mock 数据并运行全量数据管线（data/ 目录为空）", "level": 1},
        {"text": "② 训练 SASRec 基线模型 → checkpoints/sasrec.pth", "level": 1},
        {"text": "③ 训练 minGPT 生成式模型 → checkpoints/mingpt.pth", "level": 1},
        {"text": "④ 运行离线评估，验证 HR@10 提升 ≥5%（验收标准）", "level": 1},
        {"text": "⑤ 修复 constrained_beam_search 的多样性 bug（每个 beam 共享 top-k）", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ P1 应该完成（提升完成度）", "bold": True},
        {"text": "⑥ 扩充理由模板库（当前仅 5 个类目对，目标覆盖 80%+）", "level": 1},
        {"text": "⑦ 增加 TensorBoard / Wandb 训练监控", "level": 1},
        {"text": "⑧ 增加基于验证集 HR@10 的早停策略", "level": 1},
        {"text": "⑨ 将评估样本上限从 100 条提升至完整测试集", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ P2 建议完成（精进方向）", "bold": True},
        {"text": "⑩ 前端增加三模型（SASRec / minGPT / minGPT+约束）对比视图", "level": 1},
        {"text": "⑪ 增加训练进度展示与评估指标仪表盘", "level": 1},
        {"text": "⑫ Docker 部署 + YAML 配置管理 + 分层日志系统", "level": 1},
    ]
    add_content_slide(
        prs,
        "二、待完成任务清单",
        bullets2,
        "优先级排序：P0 > P1 > P2"
    )

    # ─── Slide 3: 可精进方向 ───────────────────────────────────────────────
    add_title_slide(prs, "三、可精进方向与四周开发建议", "模型 | 数据 | 理由 | 评估 | 工程")

    bullets3 = [
        {"text": "■ 模型层精进", "bold": True},
        {"text": "• minGPT：商品类别感知位置编码 + 类别条件解码 + InfoNCE 对比损失", "level": 1},
        {"text": "• SASRec：加入类别特征 Embedding、更长序列（50）、LayerNorm/LayerDrop", "level": 1},
        {"text": "• 解码器：独立 logits top-k、多样化束搜索（diverse beam search）", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ 数据层精进", "bold": True},
        {"text": "• 接入真实天猫数据集（天池）或 Amazon Electronics 子集", "level": 1},
        {"text": "• 引入 ProductKB 知识图谱辅助互补对标注", "level": 1},
        {"text": "• 序列增强（随机截断）+ popularity-based 对抗负采样", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ 理由生成精进", "bold": True},
        {"text": "• 模板库扩充（目标 80%+ 覆盖）+ TF-IDF 关键词自动抽取", "level": 1},
        {"text": "• 短/中/长三档理由模板 + 离线调用 Qwen/ChatGLM 润色改写", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ 四周开发建议", "bold": True},
        {"text": "• 第1周：数据管线 + SASRec 基线训练", "level": 1},
        {"text": "• 第2周：minGPT 训练 + beam search bug 修复", "level": 1},
        {"text": "• 第3周：端到端推理 + 理由扩充 + FastAPI 完善", "level": 1},
        {"text": "• 第4周：完整评估 + 前端仪表盘 + 答辩 PPT", "level": 1},
    ]
    add_content_slide(
        prs,
        "三、可精进方向与四周开发建议",
        bullets3,
        "核心验收标准：minGPT 生成式推荐 HR@10 相比 SASRec 基线提升 ≥5%"
    )

    prs.save(OUTPUT_PATH)
    print(f"Saved PPTX with {len(prs.slides)} slides → {OUTPUT_PATH}")

if __name__ == "__main__":
    main()
