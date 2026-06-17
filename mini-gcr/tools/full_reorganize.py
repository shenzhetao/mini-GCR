"""
Full reorganization:
1. Load the original 13-slide PPTX
2. Apply dark background (#2C2C2C) to the last 2 slides (appendix slides)
3. Add 4 progress analysis slides
4. Move the last 2 original slides (now dark) to BEFORE the QA slide
5. Save

Slide structure of original 13-slide PPTX:
  [1] 封面
  [2] 目录使用
  [3] 个性化推荐概述: 协同过滤
  [4] 个性化推荐概述:协同过滤算法
  [5] 问题定义与研究背景
  [6] 系统架构与技术选型
  [7] 数据处理
  [8] 项目计划与时间安排
  [9] 测试方法与实验应用
  [10] (empty or misc)
  [11] Q & A
  [12] 一、核心模块介绍  ← last-2 (needs dark bg, move before QA)
  [13] 二、实验结果与展望  ← last-1 (needs dark bg, move before QA)

Target final order (17 slides):
  [1-11] original slides 1-11
  [12] 一、核心模块介绍 (dark bg, was slide 12)
  [13] 二、实验结果与展望 (dark bg, was slide 13)
  [14] Q & A (QA slide)
  [15] mini-GCR 项目进度分析报告 (title slide)
  [16] 一、已完成模块分析 (table slide)
  [17] 二、待完成任务清单 (content slide)
  [18] 三、可精进方向与四周开发建议 (content slide)

Usage: python full_reorganize.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy

PPTX_PATH = r"c:/project/mini-GCR/第2次作业-12-基于生成式推荐的电商支付界面与购物车互补品推荐算法研究.pptx"
OUTPUT_PATH = r"c:/project/mini-GCR/第2次作业-12-基于生成式推荐的电商支付界面与购物车互补品推荐算法研究_final.pptx"

# Dark background matching slides 2-10
DARK_BG = RGBColor(0x2C, 0x2C, 0x2C)
TITLE_COLOR = RGBColor(0x1F, 0x4E, 0x79)
ACCENT_COLOR = RGBColor(0x2E, 0x75, 0xB6)
GRAY = RGBColor(0x59, 0x59, 0x59)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)


def set_solid_background(slide, rgb_color):
    """Set slide background to solid color, clearing any picture fill."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

    # Remove any picture fill element
    sp_tree = slide._element
    existing_bg = sp_tree.find(qn('p:bg'))
    if existing_bg is not None:
        sp_tree.remove(existing_bg)

    # Clear layout background so it doesn't override
    try:
        layout = slide.slide_layout
        lbg = layout.background
        lfill = lbg.fill
        lfill.background()
    except:
        pass


def add_title_slide(prs, title_text, subtitle_text=""):
    """Add a title/cover slide matching the dark theme."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    set_solid_background(slide, DARK_BG)

    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_COLOR
    bar.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(12.5), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle_text:
        subBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.45), Inches(12), Inches(0.6))
        stf = subBox.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(16)
        sp.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)
        sp.font.italic = True
        sp.alignment = PP_ALIGN.LEFT

    # Bottom accent line
    line = slide.shapes.add_shape(1, Inches(0), Inches(7.0), prs.slide_width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    return slide


def add_content_slide(prs, title, bullets, footer_text=""):
    """Add a bullet-point content slide matching dark theme."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_solid_background(slide, DARK_BG)

    # Header bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_COLOR
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.18), Inches(12.5), Inches(0.75))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area
    txBox2 = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        indent_level = bullet.get("level", 0)
        p.text = bullet.get("text", "")
        p.level = indent_level

        run = p.runs[0] if p.runs else p.add_run()
        run.text = p.text
        run.font.size = Pt(14.5 - indent_level * 1.5)
        run.font.color.rgb = GRAY if indent_level > 0 else RGBColor(0xCC, 0xCC, 0xCC)

        if bullet.get("bold"):
            run.font.bold = True
            run.font.color.rgb = ACCENT_COLOR

        p.space_before = Pt(3)
        p.space_after = Pt(2)

    if footer_text:
        ftBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(12), Inches(0.4))
        ftf = ftBox.text_frame
        fp = ftf.paragraphs[0]
        fp.text = footer_text
        fp.font.size = Pt(9)
        fp.font.color.rgb = GRAY
        fp.font.italic = True

    return slide


def add_table_slide(prs, title, headers, rows, footer_text=""):
    """Add a table slide matching dark theme."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_solid_background(slide, DARK_BG)

    # Header bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_COLOR
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.18), Inches(12.5), Inches(0.75))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = prs.slide_width - Inches(0.5)
    col_width = table_width / num_cols

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.25), Inches(1.2),
        table_width, Inches(num_rows * 0.42)
    ).table

    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(12)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x38, 0x38, 0x38)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            para.alignment = PP_ALIGN.LEFT

    if footer_text:
        ftBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(12), Inches(0.4))
        ftf = ftBox.text_frame
        fp = ftf.paragraphs[0]
        fp.text = footer_text
        fp.font.size = Pt(9)
        fp.font.color.rgb = GRAY
        fp.font.italic = True

    return slide


def get_slide_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 200:
                return t[:80]
    return "(无标题)"


def main():
    prs = Presentation(PPTX_PATH)
    print(f"Loaded PPTX with {len(prs.slides)} slides")

    print("\n=== Current order ===")
    for i, slide in enumerate(prs.slides):
        print(f"  [{i+1}] {get_slide_title(slide)}")

    # ── Step 1: Apply dark background to the last 2 slides (indices 11, 12) ──
    print("\n=== Step 1: Apply dark backgrounds to last 2 slides ===")
    for idx in [11, 12]:
        slide = prs.slides[idx]
        title = get_slide_title(slide)
        set_solid_background(slide, DARK_BG)
        print(f"  Slide {idx+1} '{title}': bg = #2C2C2C [OK]")

    # ── Step 2: Reorder - move last 2 slides (indices 11, 12) to before QA (index 10) ──
    print("\n=== Step 2: Move last 2 slides before QA slide ===")
    sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
    sldIds = list(sldIdLst)

    # QA slide is at index 10
    # Last 2 slides are at indices 11, 12
    sldId_qa = sldIds[10]
    sldId_last1 = sldIds[11]  # 一、核心模块介绍
    sldId_last2 = sldIds[12]  # 二、实验结果与展望

    # Remove from current positions
    sldIdLst.remove(sldId_last1)
    sldIdLst.remove(sldId_last2)

    # Insert before QA
    idx_qa = list(sldIdLst).index(sldId_qa)
    sldIdLst.insert(idx_qa, sldId_last1)
    sldIdLst.insert(idx_qa + 1, sldId_last2)

    print(f"  Moved slides 12,13 to before slide 11 (QA)")

    # ── Step 3: Add 4 new progress analysis slides ──
    print("\n=== Step 3: Adding 4 progress analysis slides ===")

    # Slide 15: Title slide
    add_title_slide(
        prs,
        "mini-GCR 项目进度分析报告",
        "生成式互补品推荐系统  |  分析日期：2026-05-23"
    )
    print(f"  [15] mini-GCR 项目进度分析报告 [OK]")

    # Slide 16: Module table
    headers = ["模块", "文件", "状态", "完成度", "说明"]
    rows = [
        ["数据管线", "generate_mock_data.py / preprocess.py", "✅ 代码完成", "80%", "数据目录为空，需先运行流水线"],
        ["SASRec 基线", "models/sasrec/model.py + train_sasrec.py", "✅ 代码完成", "90%", "未训练，checkpoints 为空"],
        ["minGPT 生成模型", "models/mingpt/model.py + decode.py", "✅ 代码完成", "85%", "存在 beam search 多样性 bug"],
        ["推理服务", "services/inference.py + reason.py", "✅ 代码完成", "85%", "静默降级，理由模板仅 5 个类目对"],
        ["评估脚本", "scripts/eval.py", "✅ 代码完成", "90%", "硬编码评测上限 100 条，未执行"],
        ["FastAPI", "app.py", "✅ 代码完成", "90%", "4 个接口正常，含 CORS"],
        ["Vue 前端", "frontend/src/App.vue", "✅ 基础完成", "70%", "缺少对比视图、指标仪表盘、商品图片"],
        ["基础设施", "config.py + run_pipeline.py", "✅ 完整", "95%", "配置中心、流水线脚本、依赖管理"],
    ]
    add_table_slide(
        prs,
        "一、已完成模块分析",
        headers, rows,
        "注：整体代码完成度约 80%，但数据目录为空、模型未训练，端到端链路尚未打通。"
    )
    print(f"  [16] 一、已完成模块分析 [OK]")

    # Slide 17: Task list
    bullets_tasks = [
        {"text": "■ P0 必须完成（影响核心交付）", "bold": True},
        {"text": "① 生成 Mock 数据并运行全量数据管线（data/ 目录为空）", "level": 1},
        {"text": "② 训练 SASRec 基线模型 → checkpoints/sasrec.pth", "level": 1},
        {"text": "③ 训练 minGPT 生成式模型 → checkpoints/mingpt.pth", "level": 1},
        {"text": "④ 运行离线评估，验证 HR@10 提升 ≥5%（验收标准）", "level": 1},
        {"text": "⑤ 修复 constrained_beam_search 的多样性 bug", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ P1 应该完成（提升完成度）", "bold": True},
        {"text": "⑥ 扩充理由模板库（当前仅 5 个类目对，目标覆盖 80%+）", "level": 1},
        {"text": "⑦ 增加 TensorBoard / Wandb 训练监控", "level": 1},
        {"text": "⑧ 增加基于验证集 HR@10 的早停策略", "level": 1},
        {"text": "⑨ 将评估样本上限从 100 条提升至完整测试集", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ P2 建议完成（精进方向）", "bold": True},
        {"text": "⑩ 前端增加三模型对比视图（SASRec / minGPT / minGPT+约束）", "level": 1},
        {"text": "⑪ 增加训练进度展示与评估指标仪表盘", "level": 1},
        {"text": "⑫ Docker 部署 + YAML 配置管理 + 分层日志系统", "level": 1},
    ]
    add_content_slide(prs, "二、待完成任务清单", bullets_tasks, "优先级排序：P0 > P1 > P2")
    print(f"  [17] 二、待完成任务清单 [OK]")

    # Slide 18: Improvement directions
    bullets_improve = [
        {"text": "■ 模型层精进", "bold": True},
        {"text": "• minGPT：商品类别感知位置编码 + 类别条件解码 + InfoNCE 对比损失", "level": 1},
        {"text": "• SASRec：类别特征 Embedding、更长序列（50）、LayerNorm/LayerDrop", "level": 1},
        {"text": "• 解码器：独立 logits top-k、多样化束搜索（diverse beam search）", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ 数据层精进", "bold": True},
        {"text": "• 接入真实天猫数据集（天池）或 Amazon Electronics 子集", "level": 1},
        {"text": "• 引入 ProductKB 知识图谱辅助互补对标注", "level": 1},
        {"text": "• 序列增强（随机截断）+ popularity-based 对抗负采样", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ 理由生成精进", "bold": True},
        {"text": "• 模板库扩充（目标 80%+ 覆盖）+ TF-IDF 关键词自动抽取", "level": 1},
        {"text": "• 短/中/长三档理由模板 + 离线调用 Qwen/ChatGLM 润色改写", "level": 1},
        {"text": " ", "bold": False},
        {"text": "■ 四周开发建议", "bold": True},
        {"text": "• 第1周：数据管线 + SASRec 基线训练", "level": 1},
        {"text": "• 第2周：minGPT 训练 + beam search bug 修复", "level": 1},
        {"text": "• 第3周：端到端推理 + 理由扩充 + FastAPI 完善", "level": 1},
        {"text": "• 第4周：完整评估 + 前端仪表盘 + 答辩 PPT", "level": 1},
    ]
    add_content_slide(
        prs,
        "三、可精进方向与四周开发建议",
        bullets_improve,
        "核心验收标准：minGPT 生成式推荐 HR@10 相比 SASRec 基线提升 ≥5%"
    )
    print(f"  [18] 三、可精进方向与四周开发建议 [OK]")

    # ── Save ──
    prs.save(OUTPUT_PATH)
    print(f"\n=== Saved to: {OUTPUT_PATH} ===")

    # Verify
    prs2 = Presentation(OUTPUT_PATH)
    print(f"\n=== Final order ({len(prs2.slides)} slides) ===")
    for i, slide in enumerate(prs2.slides):
        print(f"  [{i+1}] {get_slide_title(slide)}")


if __name__ == "__main__":
    main()
